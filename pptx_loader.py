from __future__ import annotations

from typing import List

from PIL import Image
from pptx import Presentation

from util_io import SlideContent


def _extract_text_from_shape(shape) -> str:
    parts: List[str] = []
    if hasattr(shape, "text"):
        t = (shape.text or "").strip()
        if t:
            parts.append(t)
    if hasattr(shape, "table") and shape.has_table:
        table = shape.table
        for row in table.rows:
            row_vals = []
            for cell in row.cells:
                row_vals.append((cell.text or "").strip())
            parts.append(" | ".join(v for v in row_vals if v))
    return "\n".join(p for p in parts if p)


def _extract_images_from_slide(slide) -> List[Image.Image]:
    images: List[Image.Image] = []
    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            try:
                image_blob = shape.image.blob
                img = Image.open(io.BytesIO(image_blob)).convert("RGB")
                images.append(img)
            except Exception:
                pass
    # As fallback, render slide.thumbnail if available via python-pptx? Not directly supported; omit.
    return images


import io


def load_pptx(path: str) -> List[SlideContent]:
    prs = Presentation(path)
    slides: List[SlideContent] = []
    for idx, slide in enumerate(prs.slides, start=1):
        text_parts: List[str] = []
        # Extract title if present
        title = None
        if hasattr(slide, 'shapes'):
            for shape in slide.shapes:
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    if shape.placeholder_format.type == 0:  # TITLE
                        t = (shape.text or '').strip()
                        if t:
                            title = t
                            break
        if title:
            text_parts.append(f"[TITLE] {title}")
        for shape in slide.shapes:
            try:
                t = _extract_text_from_shape(shape)
                if t:
                    text_parts.append(t)
            except Exception:
                continue
        images = _extract_images_from_slide(slide)
        slides.append(SlideContent(slide_index=idx, text="\n".join(text_parts), images=images))
    return slides


